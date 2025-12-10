import streamlit as st
import uuid
import os
import requests
import tempfile
import json
from pathlib import Path
from datetime import datetime
import io
import pickle
import subprocess
import sys

# ----------------------
# Auto-install packages
# ----------------------
def install_package(pkg):
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", pkg, "-q"])
        return True
    except:
        return False

PACKAGES = {
    'gtts': 'gtts', 'PyPDF2': 'PyPDF2', 'docx': 'python-docx',
    'pptx': 'python-pptx', 'pandas': 'pandas', 'PIL': 'Pillow',
    'deep_translator': 'deep-translator', 'fpdf': 'fpdf',
    'speech_recognition': 'SpeechRecognition',
    'audio_recorder_streamlit': 'audio-recorder-streamlit'
}

for mod, pkg in PACKAGES.items():
    try:
        __import__(mod)
    except ImportError:
        install_package(pkg)

# Imports
try:
    from gtts import gTTS
    TTS_OK = True
except:
    TTS_OK = False

try:
    import PyPDF2
    PDF_OK = True
except:
    PDF_OK = False

try:
    from docx import Document
    DOCX_OK = True
except:
    DOCX_OK = False

try:
    import pptx
    PPTX_OK = True
except:
    PPTX_OK = False

try:
    from PIL import Image
    PIL_OK = True
except:
    PIL_OK = False

try:
    from deep_translator import GoogleTranslator
    TRANS_OK = True
except:
    TRANS_OK = False

try:
    from fpdf import FPDF
    FPDF_OK = True
except:
    FPDF_OK = False

try:
    import speech_recognition as sr
    SR_OK = True
except:
    SR_OK = False

try:
    from audio_recorder_streamlit import audio_recorder
    AR_OK = True
except:
    AR_OK = False

# ----------------------
# Page Config
# ----------------------
st.set_page_config(
    page_title="Shiva AI",
    layout="wide",
    page_icon="🔱",
    initial_sidebar_state="collapsed"
)

# ----------------------
# 100% WHITE UI CSS
# ----------------------
st.markdown("""
<style>
    :root {
        color-scheme: light only !important;
    }
    
    html, body, *, *::before, *::after {
        color-scheme: light only !important;
    }
    
    html, body, div, span, p, a, li, ul, ol, h1, h2, h3, h4, h5, h6, 
    label, input, textarea, button, select, option, table, tr, td, th,
    article, section, header, footer, main, nav,
    .stApp, .main, [data-testid="stAppViewContainer"], 
    [data-testid="stAppViewContainer"] > div,
    [data-testid="stVerticalBlock"], .block-container,
    .element-container, .stMarkdown, .stText {
        background-color: #ffffff !important;
        color: #000000 !important;
    }
    
    header, [data-testid="stHeader"], [data-testid="stToolbar"] {
        background-color: #ffffff !important;
    }
    
    [data-testid="stSidebar"], [data-testid="stSidebar"] > div,
    [data-testid="stSidebar"] * {
        background-color: #f5f5f5 !important;
        color: #000000 !important;
    }
    
    /* File Uploader - WHITE */
    .stFileUploader, [data-testid="stFileUploader"],
    [data-testid="stFileUploader"] > div,
    [data-testid="stFileUploader"] section,
    [data-testid="stFileUploader"] label,
    [data-testid="stFileUploader"] span,
    [data-testid="stFileUploader"] button,
    [data-testid="stFileUploaderDropzone"],
    [data-testid="stFileUploaderDropzoneInstructions"],
    [data-testid="stFileUploaderDropzoneInstructions"] * {
        background-color: #ffffff !important;
        background: #ffffff !important;
        color: #000000 !important;
    }
    
    [data-testid="stFileUploader"] > section,
    .stFileUploader section {
        background-color: #ffffff !important;
        border: 2px dashed #cccccc !important;
        border-radius: 10px !important;
    }
    
    /* Buttons */
    .stButton > button, .stDownloadButton > button {
        background-color: #ffffff !important;
        color: #000000 !important;
        border: 2px solid #cccccc !important;
        border-radius: 10px !important;
        min-height: 45px !important;
    }
    
    .stButton > button:hover, .stDownloadButton > button:hover {
        background-color: #f0f0f0 !important;
    }
    
    /* Inputs */
    input, textarea, .stTextInput input, .stTextArea textarea,
    [data-testid="stTextInput"] input, [data-testid="stTextArea"] textarea {
        background-color: #ffffff !important;
        color: #000000 !important;
        border: 2px solid #e0e0e0 !important;
        border-radius: 10px !important;
    }
    
    /* Chat */
    [data-testid="stChatInput"], [data-testid="stChatInput"] textarea,
    .stChatInputContainer, .stChatInputContainer textarea {
        background-color: #ffffff !important;
        color: #000000 !important;
    }
    
    .stChatMessage, [data-testid="stChatMessage"] {
        background-color: #f9f9f9 !important;
        color: #000000 !important;
        border: 1px solid #e0e0e0 !important;
        border-radius: 15px !important;
    }
    
    .stChatMessage *, [data-testid="stChatMessage"] * {
        color: #000000 !important;
    }
    
    /* Select */
    .stSelectbox, .stSelectbox > div, [data-baseweb="select"],
    [data-baseweb="menu"], [data-baseweb="menu"] li, [role="option"] {
        background-color: #ffffff !important;
        color: #000000 !important;
    }
    
    [data-baseweb="menu"] li:hover, [role="option"]:hover {
        background-color: #e3f2fd !important;
    }
    
    /* Expander */
    .streamlit-expanderHeader, [data-testid="stExpander"],
    [data-testid="stExpander"] *, details, summary {
        background-color: #ffffff !important;
        color: #000000 !important;
    }
    
    /* Slider, Checkbox */
    .stSlider, .stSlider *, .stCheckbox, .stCheckbox * {
        background-color: #ffffff !important;
        color: #000000 !important;
    }
    
    /* Alerts */
    .stAlert, [data-testid="stAlert"], .stAlert * {
        background-color: #ffffff !important;
        color: #000000 !important;
    }
    
    /* Code */
    code, pre {
        background-color: #f5f5f5 !important;
        color: #000000 !important;
    }
    
    /* Custom */
    .model-box {
        border-radius: 15px;
        padding: 15px;
        text-align: center;
        margin: 10px 0;
    }
    .model-box.shiva01 {
        background: linear-gradient(135deg, #ff9933, #ff6600) !important;
    }
    .model-box.shiva02 {
        background: linear-gradient(135deg, #2196F3, #1565C0) !important;
    }
    .model-box p {
        color: #ffffff !important;
        background: transparent !important;
        margin: 0;
    }
    .model-name { font-size: 18px; font-weight: bold; }
    .model-desc { font-size: 12px; opacity: 0.9; }
    
    .chat-badge {
        display: inline-block;
        padding: 6px 15px;
        border-radius: 20px;
        font-size: 14px;
        font-weight: 600;
        margin-bottom: 10px;
    }
    .chat-badge.b01 {
        background: linear-gradient(135deg, #ff9933, #ff6600) !important;
        color: #ffffff !important;
    }
    .chat-badge.b02 {
        background: linear-gradient(135deg, #2196F3, #1565C0) !important;
        color: #ffffff !important;
    }
    
    .welcome-box {
        display: flex;
        justify-content: center;
        align-items: center;
        min-height: 50vh;
        text-align: center;
    }
    .welcome-title {
        font-size: 1.8rem;
        font-weight: 300;
        color: #000000 !important;
    }
    
    hr { border-color: #e0e0e0 !important; }
    
    @media (max-width: 768px) {
        .main .block-container { padding: 1rem 0.5rem !important; }
        .stButton > button { min-height: 50px !important; }
        h1 { font-size: 1.4rem !important; }
        
        [data-testid="stFileUploader"], [data-testid="stFileUploader"] * {
            background-color: #ffffff !important;
            background: #ffffff !important;
            color: #000000 !important;
        }
    }
    
    @media (prefers-color-scheme: dark) {
        *, html, body, .stApp, [data-testid="stFileUploader"], [data-testid="stFileUploader"] * {
            background-color: #ffffff !important;
            color: #000000 !important;
            color-scheme: light only !important;
        }
    }
</style>
""", unsafe_allow_html=True)

# ----------------------
# API Config
# ----------------------
# Shiva0.1 = Sarvam
SARVAM_KEY = os.getenv("CHATBOT_API_KEY", "sk_h4gsam68_z6e2xo8u9aaleUaBshwjVyDk")
SARVAM_URL = "https://api.sarvam.ai/v1/chat/completions"

# Shiva0.2 = Groq
GROQ_KEY = "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"
GROQ_URL = "https://api.groq.com/openai/v1/chat/completions"

# Presentation
PRES_KEY = "sk_u0xsaah7_pxzQfsbxtI4S7SsXlLLIsjaa"

# Pexels API for images
PEXELS_API_KEY = "3Y3jiJZ6WAL49N6lPsdlRbRZ6IZBfHZFHP86dr9yZfxFYoxedLLlDKAC"

SYSTEM = "You are Shiva AI, created by Shivansh Mahajan. Be helpful and professional."
LANGS = {'en': 'English', 'de': 'German', 'fr': 'French', 'es': 'Spanish'}
STORAGE = Path("shiva_data")
STORAGE.mkdir(exist_ok=True)


# ----------------------
# Pexels Image Functions
# ----------------------
def search_pexels_image(query):
    """Search for an image on Pexels and return the image URL."""
    try:
        url = "https://api.pexels.com/v1/search"
        headers = {"Authorization": PEXELS_API_KEY}
        params = {
            "query": query,
            "per_page": 1,
            "orientation": "landscape",
            "size": "large"
        }
        
        response = requests.get(url, headers=headers, params=params, timeout=10)
        
        if response.status_code == 200:
            data = response.json()
            if data.get("photos") and len(data["photos"]) > 0:
                return data["photos"][0]["src"].get("large") or data["photos"][0]["src"].get("medium")
        return None
    except Exception:
        return None


def download_image(url):
    """Download image from URL and return as BytesIO."""
    try:
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.get(url, headers=headers, timeout=15)
        
        if response.status_code == 200:
            img_data = io.BytesIO(response.content)
            
            if PIL_OK:
                from PIL import Image as PILImage
                img = PILImage.open(img_data)
                
                if img.mode in ('RGBA', 'P'):
                    img = img.convert('RGB')
                
                output = io.BytesIO()
                img.save(output, format='JPEG', quality=90)
                output.seek(0)
                return output
            else:
                img_data.seek(0)
                return img_data
        return None
    except Exception:
        return None


def generate_image_keywords(title):
    """Generate search keywords from slide title."""
    stop_words = {
        'the', 'a', 'an', 'and', 'or', 'but', 'in', 'on', 'at', 'to', 'for',
        'of', 'with', 'by', 'from', 'is', 'are', 'was', 'were', 'be', 'been',
        'being', 'have', 'has', 'had', 'do', 'does', 'did', 'will', 'would',
        'could', 'should', 'may', 'might', 'must', 'shall', 'can', 'need',
        'introduction', 'conclusion', 'overview', 'summary', 'key', 'points'
    }
    
    words = title.lower().split()
    keywords = [w for w in words if w not in stop_words and len(w) > 2]
    
    if keywords:
        return ' '.join(keywords[:3])
    return title


# ----------------------
# AI Functions
# ----------------------
def transcribe(audio):
    if not SR_OK:
        return None, "Not available"
    try:
        r = sr.Recognizer()
        with tempfile.NamedTemporaryFile(delete=False, suffix=".wav") as f:
            f.write(audio)
            p = f.name
        with sr.AudioFile(p) as s:
            a = r.record(s)
        try:
            t = r.recognize_google(a)
            os.unlink(p)
            return t, None
        except:
            os.unlink(p)
            return None, "Could not understand"
    except Exception as e:
        return None, str(e)


def ai_shiva01(msgs):
    """Shiva0.1 = Sarvam"""
    try:
        r = requests.post(SARVAM_URL, headers={"Authorization": f"Bearer {SARVAM_KEY}", "Content-Type": "application/json"},
                         json={"model": "sarvam-m", "messages": msgs, "max_tokens": 4096}, timeout=120)
        if r.status_code == 200:
            return r.json()["choices"][0]["message"]["content"]
        return f"Error {r.status_code}"
    except Exception as e:
        return f"Error: {e}"


def ai_shiva02(msgs):
    """Shiva0.2 = Groq"""
    try:
        r = requests.post(GROQ_URL, headers={"Authorization": f"Bearer {GROQ_KEY}", "Content-Type": "application/json"},
                         json={"model": "llama-3.1-8b-instant", "messages": msgs, "max_tokens": 4096}, timeout=60)
        if r.status_code == 200:
            return r.json()["choices"][0]["message"]["content"]
        return f"Error {r.status_code}"
    except Exception as e:
        return f"Error: {e}"


def ai_response(msgs, model):
    return ai_shiva01(msgs) if model == "shiva01" else ai_shiva02(msgs)


def make_flashcards(text, n=10):
    prompt = f"Generate {n} flashcards. Return ONLY JSON: [{{\"front\":\"Q\",\"back\":\"A\"}}]\n\nText: {text}"
    try:
        r = requests.post(GROQ_URL, headers={"Authorization": f"Bearer {GROQ_KEY}", "Content-Type": "application/json"},
                         json={"model": "llama-3.1-8b-instant", "messages": [{"role": "user", "content": prompt}], "max_tokens": 4000}, timeout=120)
        c = r.json()["choices"][0]["message"]["content"]
        if "```" in c:
            c = c.split("```")[1].replace("json", "")
        s, e = c.find("["), c.rfind("]") + 1
        return json.loads(c[s:e]) if s >= 0 else []
    except:
        return []


def cards_txt(cards):
    return "FLASHCARDS - Shiva AI\n" + "="*30 + "\n\n" + "\n\n".join([f"Q: {c['front']}\nA: {c['back']}" for c in cards])


def cards_word(cards):
    if not DOCX_OK:
        return None
    try:
        d = Document()
        d.add_heading("Flashcards - Shiva AI", 0)
        for i, c in enumerate(cards, 1):
            d.add_paragraph(f"{i}. Q: {c['front']}")
            d.add_paragraph(f"   A: {c['back']}")
        b = io.BytesIO()
        d.save(b)
        b.seek(0)
        return b
    except:
        return None


def cards_pdf(cards):
    if not FPDF_OK:
        return None
    try:
        p = FPDF()
        p.add_page()
        p.set_font('Arial', 'B', 14)
        p.cell(0, 10, 'Flashcards - Shiva AI', 0, 1, 'C')
        p.set_font('Arial', '', 11)
        for i, c in enumerate(cards, 1):
            if p.get_y() > 250:
                p.add_page()
            q = c['front'].encode('latin-1', 'replace').decode('latin-1')
            a = c['back'].encode('latin-1', 'replace').decode('latin-1')
            p.multi_cell(0, 6, f"{i}. Q: {q}")
            p.multi_cell(0, 6, f"   A: {a}")
            p.ln(3)
        b = io.BytesIO()
        b.write(p.output(dest='S').encode('latin-1'))
        b.seek(0)
        return b
    except:
        return None


def make_pres(topic, n):
    """Generate presentation content using Shiva0.1"""
    prompt = f"""You are an expert. Generate {n} slides for '{topic}'.
Return ONLY a JSON array: [{{"title":"...", "content":"bullet1\\nbullet2\\nbullet3"}}]
No markdown, no extra text."""
    try:
        r = requests.post(SARVAM_URL, headers={"Authorization": f"Bearer {PRES_KEY}", "Content-Type": "application/json"},
                         json={"model": "sarvam-m", "messages": [{"role": "user", "content": prompt}], "max_tokens": 4096}, timeout=120)
        c = r.json()["choices"][0]["message"]["content"]
        if "```" in c:
            c = c.split("```")[1].replace("json", "")
        s, e = c.find("["), c.rfind("]") + 1
        return json.loads(c[s:e]) if s >= 0 else []
    except:
        return []


def translate(t, lang):
    if not t or lang == 'en' or not TRANS_OK:
        return t
    try:
        return GoogleTranslator(source='en', target=lang).translate(t)
    except:
        return t


def make_pptx(slides, topic, lang_name, include_images=True, progress_callback=None):
    """Create PowerPoint with Pexels images"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ===== TITLE SLIDE =====
    if progress_callback:
        progress_callback("Creating title slide...")
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Add background image
    if include_images:
        if progress_callback:
            progress_callback("Finding title image...")
        img_url = search_pexels_image(topic)
        if img_url:
            img_data = download_image(img_url)
            if img_data:
                try:
                    slide.shapes.add_picture(
                        img_data, 
                        Inches(0), Inches(0),
                        width=prs.slide_width, 
                        height=prs.slide_height
                    )
                except:
                    pass
    
    # Dark overlay
    overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    overlay.line.fill.background()
    
    # Set transparency
    try:
        from pptx.oxml.ns import qn
        from pptx.oxml import parse_xml
        spPr = overlay._sp.spPr
        fill = spPr.find(qn('a:solidFill'))
        if fill is not None:
            srgbClr = fill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha = parse_xml('<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="50000"/>')
                srgbClr.append(alpha)
    except:
        pass
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_para = title_frame.paragraphs[0]
    title_para.text = topic
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.8))
    sub_para = sub_box.text_frame.paragraphs[0]
    sub_para.text = f"Made by Shiva AI in {lang_name}"
    sub_para.font.size = Pt(24)
    sub_para.font.color.rgb = RGBColor(200, 200, 200)
    sub_para.alignment = PP_ALIGN.CENTER
    
    # ===== CONTENT SLIDES =====
    for idx, slide_data in enumerate(slides):
        if progress_callback:
            progress_callback(f"Creating slide {idx + 1} of {len(slides)}...")
        
        content_slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Add image on right side
        img_added = False
        if include_images:
            keywords = generate_image_keywords(slide_data['title'])
            if progress_callback:
                progress_callback(f"Finding image for: {keywords}...")
            
            img_url = search_pexels_image(keywords)
            if img_url:
                img_data = download_image(img_url)
                if img_data:
                    try:
                        content_slide.shapes.add_picture(
                            img_data,
                            Inches(7.0), Inches(0),
                            width=Inches(6.333),
                            height=prs.slide_height
                        )
                        img_added = True
                    except:
                        pass
        
        # Left accent bar
        left_bar = content_slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), prs.slide_height)
        left_bar.fill.solid()
        left_bar.fill.fore_color.rgb = RGBColor(0, 120, 215)
        left_bar.line.fill.background()
        
        # Slide title
        title_width = Inches(6.0) if img_added else Inches(12)
        title_box = content_slide.shapes.add_textbox(Inches(0.4), Inches(0.5), title_width, Inches(1))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_para = title_frame.paragraphs[0]
        title_para.text = slide_data['title']
        title_para.font.size = Pt(36)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0, 0, 0)
        
        # Underline
        underline = content_slide.shapes.add_shape(1, Inches(0.4), Inches(1.4), Inches(2), Inches(0.04))
        underline.fill.solid()
        underline.fill.fore_color.rgb = RGBColor(0, 120, 215)
        underline.line.fill.background()
        
        # Content
        content_width = Inches(6.0) if img_added else Inches(12)
        content_box = content_slide.shapes.add_textbox(Inches(0.4), Inches(1.7), content_width, Inches(5.3))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        
        content_text = slide_data.get('content', '')
        content_lines = content_text.split('\n')
        
        for i, line in enumerate(content_lines):
            line = line.strip()
            if not line:
                continue
            line = line.lstrip('•-*●→▪ ')
            
            if i == 0:
                para = content_frame.paragraphs[0]
            else:
                para = content_frame.add_paragraph()
            
            para.text = f"● {line}"
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(50, 50, 50)
            para.space_before = Pt(12)
            para.space_after = Pt(6)
    
    # ===== THANK YOU SLIDE =====
    if progress_callback:
        progress_callback("Creating thank you slide...")
    
    thank_slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background image
    if include_images:
        img_url = search_pexels_image("thank you success celebration")
        if img_url:
            img_data = download_image(img_url)
            if img_data:
                try:
                    thank_slide.shapes.add_picture(
                        img_data,
                        Inches(0), Inches(0),
                        width=prs.slide_width,
                        height=prs.slide_height
                    )
                except:
                    pass
    
    # Dark overlay
    overlay = thank_slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
    overlay.line.fill.background()
    
    try:
        from pptx.oxml.ns import qn
        from pptx.oxml import parse_xml
        spPr = overlay._sp.spPr
        fill = spPr.find(qn('a:solidFill'))
        if fill is not None:
            srgbClr = fill.find(qn('a:srgbClr'))
            if srgbClr is not None:
                alpha = parse_xml('<a:alpha xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" val="55000"/>')
                srgbClr.append(alpha)
    except:
        pass
    
    # Thank you text
    thank_box = thank_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    thank_para = thank_box.text_frame.paragraphs[0]
    thank_para.text = "Thank You!"
    thank_para.font.size = Pt(72)
    thank_para.font.bold = True
    thank_para.font.color.rgb = RGBColor(255, 255, 255)
    thank_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    q_box = thank_slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12.333), Inches(0.8))
    q_para = q_box.text_frame.paragraphs[0]
    q_para.text = "Made by Shiva AI"
    q_para.font.size = Pt(28)
    q_para.font.color.rgb = RGBColor(180, 180, 180)
    q_para.alignment = PP_ALIGN.CENTER
    
    # Save
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ----------------------
# Session Helpers
# ----------------------
def uid():
    if "uid" not in st.session_state:
        st.session_state.uid = str(uuid.uuid4())
    return st.session_state.uid


def save():
    try:
        with open(STORAGE / f"{uid()}.pkl", 'wb') as f:
            pickle.dump({k: {x: y for x, y in v.items() if x != "audio"} for k, v in st.session_state.chats.items()}, f)
    except:
        pass


def load():
    try:
        p = STORAGE / f"{uid()}.pkl"
        if p.exists():
            with open(p, 'rb') as f:
                return pickle.load(f)
    except:
        pass
    return None


def speak(t):
    if not TTS_OK:
        return None
    try:
        tts = gTTS(t[:1000], lang="en")
        f = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3")
        tts.save(f.name)
        return f.name
    except:
        return None


def new_chat():
    nid = str(uuid.uuid4())
    st.session_state.chats[nid] = {"msgs": [{"role": "system", "content": SYSTEM}], "files": [], "title": "New Chat", "created": datetime.now().strftime("%Y-%m-%d %H:%M")}
    st.session_state.cid = nid
    save()


def del_chat(sid):
    if len(st.session_state.chats) > 1:
        del st.session_state.chats[sid]
        if sid == st.session_state.cid:
            st.session_state.cid = list(st.session_state.chats.keys())[0]
        save()
        return True
    return False


def extract(data, name):
    ext = Path(name).suffix.lower()
    try:
        if ext == '.pdf' and PDF_OK:
            return "\n".join([p.extract_text() or "" for p in PyPDF2.PdfReader(io.BytesIO(data)).pages])
        elif ext == '.docx' and DOCX_OK:
            return "\n".join([p.text for p in Document(io.BytesIO(data)).paragraphs])
        elif ext in ['.txt', '.py', '.js', '.json', '.md']:
            return data.decode('utf-8', errors='ignore')
    except:
        pass
    return f"[{name}]"


# ----------------------
# Init State
# ----------------------
uid()

for k, v in [("mode", "chat"), ("ppt_st", "enter"), ("fc_st", "enter"), ("cards", []), ("model", "shiva02"), ("pfiles", [])]:
    if k not in st.session_state:
        st.session_state[k] = v

if "chats" not in st.session_state:
    ld = load()
    if ld:
        st.session_state.chats = ld
        st.session_state.cid = list(ld.keys())[0]
    else:
        st.session_state.chats = {}
        new_chat()

if "cid" not in st.session_state:
    st.session_state.cid = list(st.session_state.chats.keys())[0]


# ----------------------
# Sidebar
# ----------------------
with st.sidebar:
    st.markdown("## 🔱 Shiva AI")
    st.markdown("*By Shivansh Mahajan*")
    st.markdown("---")
    
    # Model Selection
    st.markdown("### 🧠 Model")
    
    if st.session_state.model == "shiva01":
        st.markdown('<div class="model-box shiva01"><p class="model-name">🔱 Shiva0.1</p><p class="model-desc">Sarvam - Good for Hindi</p></div>', unsafe_allow_html=True)
        if st.button("🔄 Switch to Shiva0.2", use_container_width=True):
            st.session_state.model = "shiva02"
            st.rerun()
    else:
        st.markdown('<div class="model-box shiva02"><p class="model-name">🚀 Shiva0.2</p><p class="model-desc">Fast & multilingual</p></div>', unsafe_allow_html=True)
        if st.button("🔄 Switch to Shiva0.1", use_container_width=True):
            st.session_state.model = "shiva01"
            st.rerun()
    
    st.markdown("---")
    
    if st.button("✨ New Chat", use_container_width=True):
        st.session_state.mode = "chat"
        new_chat()
        st.rerun()
    
    if st.button("📊 Presentation", use_container_width=True):
        st.session_state.mode = "pres"
        st.session_state.ppt_st = "enter"
        st.rerun()
    
    if st.button("📚 Flashcards", use_container_width=True):
        st.session_state.mode = "flash"
        st.session_state.fc_st = "enter"
        st.rerun()
    
    st.markdown("---")
    
    if st.session_state.mode == "chat":
        st.markdown("### 💬 Chats")
        for sid, s in sorted(st.session_state.chats.items(), key=lambda x: x[1].get("created", ""), reverse=True):
            c1, c2 = st.columns([5, 1])
            lbl = ("▶ " if sid == st.session_state.cid else "") + s["title"][:20]
            if c1.button(lbl, key=f"c_{sid}", use_container_width=True):
                st.session_state.cid = sid
                st.rerun()
            if c2.button("🗑️", key=f"d_{sid}"):
                del_chat(sid)
                st.rerun()


# ----------------------
# Main Content
# ----------------------

# FLASHCARDS
if st.session_state.mode == "flash":
    st.markdown("# 📚 Flashcards")
    st.markdown("*Made by Shiva AI*")
    st.markdown("---")
    
    if st.session_state.fc_st == "enter":
        txt = st.text_area("📝 Text:", height=180)
        num = st.slider("Cards", 5, 20, 10)
        
        if st.button("🚀 Generate", type="primary", use_container_width=True):
            if txt and len(txt) >= 50:
                with st.spinner("Creating..."):
                    cards = make_flashcards(txt, num)
                    if cards:
                        st.session_state.cards = cards
                        st.session_state.fc_st = "view"
                        st.rerun()
            else:
                st.error("Need 50+ characters")
        
        if st.button("🏠 Back"):
            st.session_state.mode = "chat"
            st.rerun()
    
    else:
        cards = st.session_state.cards
        st.success(f"✅ {len(cards)} cards!")
        
        for i, c in enumerate(cards, 1):
            with st.expander(f"{i}. {c.get('front', '')[:40]}..."):
                st.info(f"**Q:** {c.get('front', '')}")
                st.success(f"**A:** {c.get('back', '')}")
        
        st.markdown("### 📥 Download")
        c1, c2, c3 = st.columns(3)
        c1.download_button("📄 TXT", cards_txt(cards), "cards.txt", use_container_width=True)
        w = cards_word(cards)
        if w:
            c2.download_button("📘 WORD", w, "cards.docx", use_container_width=True)
        p = cards_pdf(cards)
        if p:
            c3.download_button("📕 PDF", p, "cards.pdf", use_container_width=True)
        
        if st.button("🔄 More"):
            st.session_state.fc_st = "enter"
            st.rerun()
        if st.button("🏠 Back"):
            st.session_state.mode = "chat"
            st.rerun()


# PRESENTATION
elif st.session_state.mode == "pres":
    st.markdown("# 📊 Presentation Generator")
    st.markdown("*Made by Shiva AI with HD Images*")
    st.markdown("---")
    
    if st.session_state.ppt_st == "enter":
        topic = st.text_input("📝 Topic:")
        lang = st.selectbox("🌐 Language", list(LANGS.keys()), format_func=lambda x: LANGS[x])
        n = st.slider("📊 Slides", 3, 10, 5)
        include_images = st.checkbox("🖼️ Include HD images from Pexels", value=True)
        
        if st.button("🚀 Create Presentation", type="primary", use_container_width=True):
            if topic:
                status_text = st.empty()
                
                def update_status(msg):
                    status_text.info(msg)
                
                with st.spinner("Creating presentation..."):
                    try:
                        update_status("🤖 Generating content...")
                        slides = make_pres(topic, n)
                        
                        if not slides:
                            st.error("Failed to generate slides. Try again.")
                        else:
                            update_status(f"🌐 Translating to {LANGS[lang]}...")
                            t_topic = translate(topic, lang)
                            t_slides = [{"title": translate(s["title"], lang), "content": translate(s["content"], lang)} for s in slides]
                            
                            update_status("📊 Building PowerPoint with images...")
                            ppt = make_pptx(t_slides, t_topic, LANGS[lang], include_images, update_status)
                            
                            st.session_state.ppt = ppt
                            st.session_state.ppt_name = f"{topic.replace(' ', '_')}_shiva_ai.pptx"
                            st.session_state.ppt_st = "done"
                            st.rerun()
                    except Exception as e:
                        st.error(f"Error: {e}")
        
        if st.button("🏠 Back"):
            st.session_state.mode = "chat"
            st.rerun()
    
    else:
        st.balloons()
        st.success("✅ Presentation ready with HD images!")
        st.download_button("📥 Download PPTX", st.session_state.ppt, st.session_state.ppt_name, use_container_width=True)
        
        if st.button("🔄 Create Another"):
            st.session_state.ppt_st = "enter"
            st.rerun()
        if st.button("🏠 Back"):
            st.session_state.mode = "chat"
            st.rerun()


# CHAT
else:
    chat = st.session_state.chats[st.session_state.cid]
    
    # Badge
    if st.session_state.model == "shiva01":
        st.markdown('<span class="chat-badge b01">🔱 Shiva0.1</span>', unsafe_allow_html=True)
    else:
        st.markdown('<span class="chat-badge b02">🚀 Shiva0.2</span>', unsafe_allow_html=True)
    
    # Messages
    user_msgs = [m for m in chat["msgs"] if m["role"] == "user"]
    
    if not user_msgs:
        st.markdown('<div class="welcome-box"><h1 class="welcome-title">🔱 How can Shiva AI help?</h1></div>', unsafe_allow_html=True)
    else:
        for m in chat["msgs"]:
            if m["role"] == "system":
                continue
            with st.chat_message(m["role"]):
                st.markdown(m["content"])
                if m["role"] == "assistant" and m.get("audio"):
                    st.audio(m["audio"])
    
    st.markdown("---")
    
    # Voice
    if AR_OK:
        st.markdown("### 🎤 Voice")
        c1, c2 = st.columns([1, 4])
        with c1:
            aud = audio_recorder(text="", recording_color="#e74c3c", neutral_color="#2196F3", icon_size="2x")
        with c2:
            if aud:
                with st.spinner("Processing..."):
                    txt, err = transcribe(aud)
                    if txt:
                        st.success(f"📝 {txt}")
                        if st.button("📤 Send"):
                            if chat["title"] == "New Chat":
                                chat["title"] = txt[:25] + "..."
                            chat["msgs"].append({"role": "user", "content": txt})
                            save()
                            msgs = [{"role": "system", "content": SYSTEM}] + [{"role": m["role"], "content": m["content"]} for m in chat["msgs"][-10:] if m["role"] in ["user", "assistant"]]
                            resp = ai_response(msgs, st.session_state.model)
                            chat["msgs"].append({"role": "assistant", "content": resp, "audio": speak(resp)})
                            save()
                            st.rerun()
                    elif err:
                        st.error(err)
        st.markdown("---")
    
    # Files
    files = st.file_uploader("📎 Files", accept_multiple_files=True, type=['pdf', 'docx', 'txt', 'py', 'json', 'md'])
    
    if files:
        names = [f.name for f in files]
        if names != st.session_state.pfiles:
            st.session_state.pfiles = names
            for f in files:
                data = f.read()
                content = extract(data, f.name)
                if "files" not in chat:
                    chat["files"] = []
                if not any(x["name"] == f.name for x in chat["files"]):
                    chat["files"].append({"name": f.name, "content": content, "size": len(data)})
            save()
            st.rerun()
    
    if chat.get("files"):
        for i, f in enumerate(chat["files"]):
            c1, c2 = st.columns([6, 1])
            c1.write(f"📎 {f['name']}")
            if c2.button("❌", key=f"rf_{i}"):
                chat["files"] = [x for x in chat["files"] if x["name"] != f["name"]]
                save()
                st.rerun()
        st.markdown("---")
    
    # Input
    inp = st.chat_input("Ask Shiva AI...")
    
    if inp:
        if chat["title"] == "New Chat":
            chat["title"] = inp[:25] + "..."
        
        chat["msgs"].append({"role": "user", "content": inp})
        save()
        
        with st.chat_message("user"):
            st.markdown(inp)
        
        ctx = ""
        if chat.get("files"):
            ctx = "\n[FILES]\n" + "\n".join([f"[{f['name']}]\n{f['content'][:2000]}" for f in chat["files"]]) + "\n[/FILES]\n\n"
        
        with st.chat_message("assistant"):
            ph = st.empty()
            nm = "Shiva0.1" if st.session_state.model == "shiva01" else "Shiva0.2"
            ph.markdown(f"🔱 {nm} thinking...")
            
            try:
                msgs = [{"role": "system", "content": SYSTEM}]
                recent = [m for m in chat["msgs"] if m["role"] in ["user", "assistant"]][-10:]
                for m in recent[:-1]:
                    msgs.append({"role": m["role"], "content": m["content"][:2000]})
                msgs.append({"role": "user", "content": ctx + inp if ctx else inp})
                resp = ai_response(msgs, st.session_state.model)
            except Exception as e:
                resp = f"Error: {e}"
            
            ph.markdown(resp)
            aud = speak(resp) if TTS_OK else None
            chat["msgs"].append({"role": "assistant", "content": resp, "audio": aud})
            save()
            if aud:
                st.audio(aud)
            st.rerun()


# Footer
st.markdown("---")
st.markdown('<div style="text-align:center;padding:15px;"><p><strong>🔱 Shiva AI</strong> - Made by Shivansh Mahajan</p></div>', unsafe_allow_html=True)
